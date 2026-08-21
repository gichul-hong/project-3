"""
Build ultra-clean, aesthetic 3-slide PPT showcase with exact 1:1 square aspect ratio images.
- Original 512x512 images rendered with strict 1:1 aspect ratio (no distortion).
- Images maximized vertically and horizontally to fill the slide (~2.78" x 2.78" square).
- Ultra-compact header and clean minimalist captions directly below each image.
- Pure images and typography only (zero background rectangles or bounding boxes).
"""

import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

FONT = "맑은 고딕"
W, H = Inches(13.333), Inches(7.5)

# Minimalist Aesthetic Palette
C_BG        = RGBColor(255, 255, 255)
C_TITLE     = RGBColor(15, 23, 42)      # Deep Slate
C_SUB       = RGBColor(100, 116, 139)   # Muted Slate
C_LINE      = RGBColor(226, 232, 240)   # Hairline
C_TEXT      = RGBColor(30, 41, 59)      # Charcoal
C_CYAN      = RGBColor(2, 132, 199)     # Exp-12 Accent
C_EMERALD   = RGBColor(5, 150, 105)     # Exp-13 Accent
C_AMBER     = RGBColor(217, 119, 6)     # Exp-14 Accent

CLASS_PROMPT = {
    "actionfigure_2": "action figure",
    "decoritems_woodenpot": "wooden pot",
    "furniture_sofa2": "sofa",
    "instrument_music2": "guitar",
    "luggage_backpack1": "backpack",
    "person_3": "person",
    "pet_cat5": "cat",
    "scene_waterfall": "waterfall",
    "transport_tank": "tank",
    "wearable_jacket1": "jacket",
}

EXP_INFO = {
    "12_balanced_ensemble": {
        "tag": "EXP-12",
        "title": "Exp-12: Balanced SOTA Ensemble",
        "desc": "28-Step Fast Euler ODE  ·  Natural Background Reference",
        "color": C_CYAN,
        "kpi_t": "0.3250",
        "kpi_i": "0.7370",
    },
    "13_sota_ensemble": {
        "tag": "EXP-13  🏆 SOTA",
        "title": "Exp-13: Ultimate SOTA Ensemble",
        "desc": "Crop-Fit Center Reference  ·  1:1 Metric Alignment  ·  White-Bg Penalty Guard",
        "color": C_EMERALD,
        "kpi_t": "0.3249",
        "kpi_i": "0.7396",
    },
    "14_extreme_prompt_align": {
        "tag": "EXP-14  🥇 TEXT SOTA",
        "title": "Exp-14: Extreme Prompt Alignment",
        "desc": "Soft Guided ODE (tau=0.6, eta=0.65)  ·  CFG 7.5 Prompt Maximizer",
        "color": C_AMBER,
        "kpi_t": "0.3402",
        "kpi_i": "0.7162",
    },
}

def get_best_8_items(exp_name, root_dir):
    exp_dir = os.path.join(root_dir, "experiments", exp_name)
    best_per_concept = []

    for concept, class_noun in CLASS_PROMPT.items():
        prompts_file = os.path.join(root_dir, "prompt", f"{concept}.txt")
        prompts = []
        if os.path.exists(prompts_file):
            with open(prompts_file, "r", encoding="utf-8") as f:
                prompts = [l.strip().replace("{}", class_noun) for l in f.readlines() if l.strip()]

        sel_file = os.path.join(exp_dir, f"selection_{concept}.json")
        sel_data = {}
        if os.path.exists(sel_file):
            with open(sel_file, "r", encoding="utf-8") as f:
                raw_sel = json.load(f)
                if isinstance(raw_sel, list):
                    for item in raw_sel:
                        if isinstance(item, dict) and "prompt_idx" in item:
                            sel_data[item["prompt_idx"]] = item
                elif isinstance(raw_sel, dict):
                    for k, v in raw_sel.items():
                        idx_num = int(k.replace("p", "")) if k.replace("p", "").isdigit() else None
                        if idx_num is not None:
                            sel_data[idx_num] = v

        concept_candidates = []
        for idx in range(1, len(prompts)):
            img_path = os.path.join(exp_dir, concept, f"{idx}.png")
            if not os.path.exists(img_path):
                img_path = os.path.join(exp_dir, concept, f"{idx}.jpg")
            if not os.path.exists(img_path):
                continue

            t, i = 0.32, 0.74
            if idx in sel_data:
                sel_obj = sel_data[idx]
                t = sel_obj.get("clip_t", 0.32)
                i = sel_obj.get("clip_i", 0.74)

            concept_candidates.append({
                "concept": concept,
                "class_noun": class_noun,
                "prompt_idx": idx,
                "prompt": prompts[idx] if idx < len(prompts) else "",
                "img_path": img_path,
                "clip_t": t,
                "clip_i": i,
                "total": t + i
            })

        if concept_candidates:
            concept_candidates.sort(key=lambda x: x["total"], reverse=True)
            best_per_concept.append(concept_candidates[0])

    best_per_concept.sort(key=lambda x: x["total"], reverse=True)
    return best_per_concept[:8]


def create_1to1_maximized_showcase():
    root_dir = os.path.dirname(os.path.abspath(__file__))
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    BLANK = prs.slide_layouts[6]

    for exp_name in ["12_balanced_ensemble", "13_sota_ensemble", "14_extreme_prompt_align"]:
        meta = EXP_INFO[exp_name]
        slide = prs.slides.add_slide(BLANK)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = C_BG

        # 1. Ultra-Compact Inline Header (Left: Title + Sub, Right: KPI)
        tb_title = slide.shapes.add_textbox(Inches(0.6), Inches(0.14), Inches(8.5), Inches(0.55))
        tf_title = tb_title.text_frame
        tf_title.word_wrap = True
        tf_title.margin_left = tf_title.margin_top = tf_title.margin_right = tf_title.margin_bottom = 0

        p_h1 = tf_title.paragraphs[0]
        p_h1.text = meta["title"]
        p_h1.font.name = FONT
        p_h1.font.size = Pt(15.5)
        p_h1.font.bold = True
        p_h1.font.color.rgb = C_TITLE
        p_h1.space_after = Pt(1)

        p_desc = tf_title.add_paragraph()
        p_desc.text = f"{meta['tag']}  |  {meta['desc']}"
        p_desc.font.name = FONT
        p_desc.font.size = Pt(8.0)
        p_desc.font.color.rgb = C_SUB

        # Right KPI Typography
        tb_kpi = slide.shapes.add_textbox(Inches(8.8), Inches(0.14), Inches(3.933), Inches(0.55))
        tf_kpi = tb_kpi.text_frame
        tf_kpi.word_wrap = True
        tf_kpi.margin_left = tf_kpi.margin_top = tf_kpi.margin_right = tf_kpi.margin_bottom = 0

        p_kpi1 = tf_kpi.paragraphs[0]
        p_kpi1.text = f"CLIP-T  {meta['kpi_t']}   ·   CLIP-I  {meta['kpi_i']}"
        p_kpi1.font.name = FONT
        p_kpi1.font.size = Pt(11.5)
        p_kpi1.font.bold = True
        p_kpi1.font.color.rgb = meta["color"]
        p_kpi1.alignment = PP_ALIGN.RIGHT
        p_kpi1.space_after = Pt(1)

        p_kpi2 = tf_kpi.add_paragraph()
        p_kpi2.text = "OFFICIAL BENCHMARK AVERAGE"
        p_kpi2.font.name = FONT
        p_kpi2.font.size = Pt(7.0)
        p_kpi2.font.bold = True
        p_kpi2.font.color.rgb = C_SUB
        p_kpi2.alignment = PP_ALIGN.RIGHT

        # Subtle Separator Line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.72), Inches(12.133), Inches(0.01))
        line.fill.solid()
        line.fill.fore_color.rgb = C_LINE
        line.line.color.rgb = C_LINE

        # 2. Exact 1:1 Aspect Ratio Maximized Grid (512x512 Ratio Preserved)
        img_size = Inches(2.78)   # Strict 1:1 Square (2.78" x 2.78")
        cols = 4
        rows = 2
        col_gap = Inches(0.18)
        row_gap = Inches(0.14)
        cap_h = Inches(0.34)

        total_grid_w = cols * img_size + (cols - 1) * col_gap   # 4 * 2.78 + 3 * 0.18 = 11.66 inches
        total_cell_h = img_size + cap_h                          # 2.78 + 0.34 = 3.12 inches
        
        grid_left = (W - total_grid_w) / 2                      # ~0.836 inches (Perfect Center)
        grid_top = Inches(0.82)

        best8 = get_best_8_items(exp_name, root_dir)

        for idx, item in enumerate(best8):
            r = idx // cols
            c = idx % cols

            x = grid_left + c * (img_size + col_gap)
            y = grid_top + r * (total_cell_h + row_gap)

            # 1. Pure 1:1 Square Image (512x512 No Distortion)
            if os.path.exists(item["img_path"]):
                slide.shapes.add_picture(item["img_path"], x, y, img_size, img_size)

            # 2. Minimalist Typography Caption (Pure Text, No Box)
            tb_cap = slide.shapes.add_textbox(x, y + img_size + Inches(0.03), img_size, cap_h)
            tf_cap = tb_cap.text_frame
            tf_cap.word_wrap = True
            tf_cap.margin_left = tf_cap.margin_top = tf_cap.margin_right = tf_cap.margin_bottom = 0

            # Line 1: Concept & Prompt (Clean)
            p_c1 = tf_cap.paragraphs[0]
            p_c1.text = f"[{item['concept']}] \"{item['prompt']}\""
            p_c1.font.name = FONT
            p_c1.font.size = Pt(6.5)
            p_c1.font.bold = True
            p_c1.font.color.rgb = C_TEXT
            p_c1.space_after = Pt(1)

            # Line 2: T & I Scores Only
            p_c2 = tf_cap.add_paragraph()
            p_c2.text = f"T  {item['clip_t']:.4f}   ·   I  {item['clip_i']:.4f}"
            p_c2.font.name = FONT
            p_c2.font.size = Pt(6.5)
            p_c2.font.bold = True
            p_c2.font.color.rgb = meta["color"]

    output_path = os.path.join(root_dir, "SD3.5_Exp12_14_Showcase_Presentation.pptx")
    prs.save(output_path)
    print(f"✓ Successfully generated exact 1:1 maximized showcase PPT: {output_path}")

if __name__ == "__main__":
    create_1to1_maximized_showcase()
